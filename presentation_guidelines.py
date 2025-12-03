from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """Create a PowerPoint presentation following the project guidelines."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1) Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Nifty 50 index Analysis using\n Data Science & Data Analytics techniques"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    # Add course title
    course_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.5))
    course_frame = course_box.text_frame
    course_frame.text = "Data Science Institute & Postgraduate Diplone in Data Analytics"
    course_para = course_frame.paragraphs[0]
    course_para.font.size = Pt(20)
    course_para.alignment = PP_ALIGN.CENTER

    # Add group members
    members_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1.5))
    members_frame = members_box.text_frame
    members_frame.text = "Group Members:\n Barry, Eoin, Finn, Ismael, Oksana, Pablo, Werner, Zuzana"
    members_para = members_frame.paragraphs[0]
    members_para.font.size = Pt(16)
    members_para.alignment = PP_ALIGN.CENTER

    # Add date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.4))
    date_frame = date_box.text_frame
    date_frame.text = "November 2024"
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(14)
    date_para.alignment = PP_ALIGN.CENTER
    date_para.font.color.rgb = RGBColor(102, 102, 102)

    # 2) Project Background - Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "Project Background"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Market Analysis Challenge"

    p = tf.add_paragraph()
    p.text = "Global financial markets are interconnected and influence each other"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Nifty 50 (India) and Dow Jones (US) are major indices representing significant economies"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Understanding cross-market relationships can help predict market movements"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Machine learning offers sophisticated tools for pattern recognition in financial data"
    p.level = 1

    # 3) Project Background - Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Why This Matters"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Practical Applications"

    p = tf.add_paragraph()
    p.text = "Investor Decision Making: Better informed trading strategies"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Risk Management: Anticipate market volatility and downturns"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Portfolio Optimization: Leverage cross-market correlations"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Academic Research: Contribute to understanding of global market dynamics"
    p.level = 1

    # 4) Data Source
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Data Source"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Yahoo Finance API"

    p = tf.add_paragraph()
    p.text = "Daily market data from January 2019 to April 2025"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Six major global indices: Nifty 50, Dow Jones, Nasdaq, Hang Seng, Nikkei 225, DAX"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "VIX (Volatility Index) as market sentiment indicator"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Twitter/X data for sentiment analysis (Nifty 50 generated posts, thanks Werner!)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Features: Open, Close prices, Returns, Year, Quarter, Month"
    p.level = 1

    # 5) Data Snapshot
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Data Snapshot"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Dataset Characteristics"

    p = tf.add_paragraph()
    p.text = "Time Period: 6+ years of daily trading data"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Total Records: ~1,500 trading days (after cleaning)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Variables: 20+ features including returns and temporal indicators"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Data Quality: Missing values handled via linear interpolation"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Target Variable: Nifty_Dir_Open (Binary: Up=1, Down=0)"
    p.level = 1

    # 6) Objectives - Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Objectives"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Primary Goals"

    p = tf.add_paragraph()
    p.text = "Perform EDA on the data"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Compare performance of multiple ML algorithms (Logistic Regression, Naive Bayes, Decision Tree, Random Forest)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Analyze correlations between global markets and Nifty 50 movements"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Incorporate sentiment analysis from social media to enhance predictions"
    p.level = 1

    # 7) Objectives - Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Success Criteria"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Measuring Success"

    p = tf.add_paragraph()
    p.text = "Model Accuracy: Achieve >60% prediction accuracy on test data"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "AUC Score: Target AUC > 0.65 for ROC curves"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Balanced Performance: Maintain similar metrics between train and test datasets"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Feature Importance: Identify key predictive variables"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Practical Insights: Generate actionable trading signals"
    p.level = 1

    # 8) Analysis Plan - Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Analysis Plan - Phase 1 & 2"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Phase 1: Data Collection & Preprocessing"

    p = tf.add_paragraph()
    p.text = "Data extraction from Yahoo Finance API"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Handle missing values and outliers"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Calculate daily returns for all indices"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Phase 2: Exploratory Data Analysis"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Correlation analysis between markets"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Time series visualization and trend analysis"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Statistical tests (normality, stationarity)"
    p.level = 1

    # 9) Analysis Plan - Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Analysis Plan - Phase 3 & 4"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Phase 3: Predictive Modeling (Logistic Regression)"

    p = tf.add_paragraph()
    p.text = "Train/test split (80/20)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "VIF analysis for multicollinearity"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "ROC curve and optimal threshold determination"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Phase 4: Advanced ML Models"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Naive Bayes Classifier"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Decision Tree with depth optimization"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Random Forest ensemble method"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Model Comparison & Final Recommendations"
    p.level = 0

    # 10) Analysis Plan - Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Analysis Plan - Phase 5"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Phase 5: Text Mining & Sentiment Analysis"

    p = tf.add_paragraph()
    p.text = "Text preprocessing (cleaning, tokenization, stopword removal)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Word cloud generation for topic identification"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "VADER sentiment analysis on Twitter/X data"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Sentiment scoring (Positive, Negative, Neutral)"
    p.level = 1


    # 11) Challenges
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Challenges Encountered"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Data Challenges"

    p = tf.add_paragraph()
    p.text = "Different trading hours and holidays across global markets"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Missing data requiring interpolation strategies"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Modeling Challenges"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Low predictive power (AUC ~0.53-0.59) indicates market complexity"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Technical Challenges"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Integration of multiple data sources and formats"
    p.level = 1

    # 12) Project Plan with Timeline
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Timeline"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.size = 4
    tf.text = "Week 1-2: Data Collection & Preprocessing"

    p = tf.add_paragraph()
    p.text = "Set up data pipeline, clean and prepare datasets"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Week 3-4: Exploratory Data Analysis"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Statistical analysis, visualization, correlation studies"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Week 5-6: Model Development"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Build and train ML models, optimize parameters"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Week 7-8: Sentiment Analysis & Integration"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "NLP processing, combine with market models"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Week 9-10: Finalization & Presentation"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Results analysis, documentation, presentation preparation"
    p.level = 1

    # Save presentation
    prs.save('Project_Presentation.pptx')
    print("Presentation created successfully: Project_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()