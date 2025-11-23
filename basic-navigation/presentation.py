"""
Generate PowerPoint presentation from Market Analysis Dashboard
Uses python-pptx to create slides with content from all phases
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
from pathlib import Path
import sys

# Add the basic-navigation directory to the path
sys.path.insert(0, str(Path(__file__).parent / 'basic-navigation'))

import phase1
import phase2
import phase3
import phase4
import phase5

# Constants
TITLE_FONT_SIZE = Pt(40)
SUBTITLE_FONT_SIZE = Pt(24)
CONTENT_FONT_SIZE = Pt(18)
DARK_BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(68, 114, 196)


def create_title_slide(prs, title, subtitle=""):
    """Create a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
    title_shape.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    title_shape.text_frame.paragraphs[0].font.bold = True

    if subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        subtitle_shape.text_frame.paragraphs[0].font.size = SUBTITLE_FONT_SIZE
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = LIGHT_BLUE

    return slide


def create_content_slide(prs, title, content_points):
    """Create a slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE

    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()

    for point in content_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = CONTENT_FONT_SIZE
        p.space_after = Pt(12)

    return slide


def create_image_slide(prs, title, image_path, caption=""):
    """Create a slide with an image"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True

    # Add image
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(5)

    try:
        slide.shapes.add_picture(str(image_path), left, top, width=width)
    except Exception as e:
        print(f"Error adding image {image_path}: {e}")

    # Add caption if provided
    if caption:
        left = Inches(1)
        top = Inches(6.8)
        width = Inches(8)
        height = Inches(0.5)

        caption_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = caption_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

    return slide


def save_plot_to_file(fig, filename):
    """Save matplotlib figure to file"""
    output_dir = Path(__file__).parent / 'temp_plots'
    output_dir.mkdir(exist_ok=True)
    filepath = output_dir / filename

    if fig is None:
        # If function didn't return fig, get current figure
        fig = plt.gcf()

    if fig is not None:
        fig.savefig(filepath, dpi=150, bbox_inches='tight')
        plt.close(fig)
    else:
        print(f"Warning: Could not save {filename} - no figure available")
        return None

    return filepath


def create_presentation():
    """Main function to create the presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title Slide
    create_title_slide(
        prs,
        "Market Analysis Dashboard",
        "Nifty50 Predictive Modeling & Sentiment Analysis"
    )

    # Table of Contents
    create_content_slide(prs, "Agenda", [
        "Phase 1: Data Collection and Preprocessing",
        "Phase 2: Exploratory Data Analysis",
        "Phase 3: Logistic Regression Model",
        "Phase 4: Machine Learning Models Comparison",
        "Phase 5: Text Mining & Sentiment Analysis",
        "Conclusions and Recommendations"
    ])

    # ========================================================================
    # PHASE 1 - Data Collection
    # ========================================================================

    create_title_slide(prs, "Phase 1", "Data Collection and Preprocessing")

    create_content_slide(prs, "Dataset Overview", [
        "6 Major Global Market Indices:",
        "  • Nifty 50 (India)",
        "  • Dow Jones (USA)",
        "  • Nasdaq (USA)",
        "  • Hang Seng (Hong Kong)",
        "  • Nikkei 225 (Japan)",
        "  • DAX (Germany)",
        "VIX Volatility Index included",
        "Time period: January 2019 to April 2025",
        "Missing values handled with linear interpolation"
    ])

    # QQ Plots
    print("Generating Phase 1 plots...")
    plt.ioff()  # Turn off interactive mode
    fig = phase1.qq_plots_for_returns()
    image_path = save_plot_to_file(fig, 'phase1_qq_plots.png')
    if image_path:
        create_image_slide(prs, "Normality Check - QQ Plots", image_path,
                          "QQ plots show daily returns approximately follow normal distribution")

    # ========================================================================
    # PHASE 2 - EDA
    # ========================================================================

    create_title_slide(prs, "Phase 2", "Exploratory Data Analysis")

    create_content_slide(prs, "Phase 2 Objectives", [
        "Analyze 5-year performance of global indices",
        "Examine correlations between markets",
        "Study impact of COVID-19 pandemic",
        "Investigate Nifty 50 daily movements",
        "Identify relationships with global markets"
    ])

    # Box-Whisker Plot
    print("Generating Phase 2 plots...")
    fig = phase2.boxplot_by_year()
    image_path = save_plot_to_file(fig, 'phase2_boxplot.png')
    if image_path:
        create_image_slide(prs, "Daily Returns Distribution by Year", image_path,
                          "Box-whisker plots showing volatility patterns across years")

    # Correlation Matrix
    fig = phase2.correlation_matrix()
    image_path = save_plot_to_file(fig, 'phase2_correlation.png')
    if image_path:
        create_image_slide(prs, "Market Correlation Analysis", image_path,
                          "Correlation matrix reveals strong relationships between markets")

    # COVID Analysis
    fig = phase2.plot_covid_recovery()
    image_path = save_plot_to_file(fig, 'phase2_covid.png')
    if image_path:
        create_image_slide(prs, "COVID-19 Impact Analysis", image_path,
                          "Market behavior during pandemic phases")

    # Nifty Direction Analysis
    fig = phase2.global_indices_against_nifty()
    image_path = save_plot_to_file(fig, 'phase2_nifty_dir.png')
    if image_path:
        create_image_slide(prs, "Global Indices vs Nifty Direction", image_path,
                          "Relationship between global markets and Nifty 50 movements")    # ========================================================================
    # PHASE 3 - Logistic Regression
    # ========================================================================

    create_title_slide(prs, "Phase 3", "Binary Logistic Regression Model")

    create_content_slide(prs, "Logistic Regression Approach", [
        "Objective: Predict Nifty 50 direction (Up/Down)",
        "Initial Model: All 6 market return variables",
        "VIF Analysis: Check for multicollinearity",
        "Refined Model: 4 significant variables",
        "  • HangSeng_Return",
        "  • Nikkei_Return",
        "  • DAX_Return",
        "  • VIX_Return",
        "Optimal threshold: 0.4753"
    ])

    # Train models first
    print("Generating Phase 3 plots...")
    phase3.logistic_regression_summary()
    _, _ = phase3.logistic_regression_significant_vars_summary()

    # ROC Curves - create manually since phase3 uses shared plotting
    from sklearn.metrics import roc_curve, auc
    import statsmodels.api as sm

    # Get predictions for ROC curves
    X_train = sm.add_constant(phase3.training_set[phase3.SIGNIFICANT_VARS])
    X_test = sm.add_constant(phase3.testing_set[phase3.SIGNIFICANT_VARS])
    y_train = phase3.training_set['Nifty_Dir_Open']
    y_test = phase3.testing_set['Nifty_Dir_Open']

    train_probs = phase3.model.predict(X_train.astype(float))
    test_probs = phase3.model.predict(X_test.astype(float))

    fpr_train, tpr_train, _ = roc_curve(y_train, train_probs)
    fpr_test, tpr_test, _ = roc_curve(y_test, test_probs)
    auc_train = auc(fpr_train, tpr_train)
    auc_test = auc(fpr_test, tpr_test)

    # Create ROC figure
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 5))

    ax1.plot(fpr_train, tpr_train, color='darkorange', lw=2, label=f'ROC (AUC = {auc_train:.2f})')
    ax1.plot([0, 1], [0, 1], color='navy', lw=2, linestyle='--')
    ax1.set_xlim([0.0, 1.0])
    ax1.set_ylim([0.0, 1.05])
    ax1.set_xlabel('False Positive Rate')
    ax1.set_ylabel('True Positive Rate')
    ax1.set_title('ROC Curve - Train Dataset')
    ax1.legend(loc="lower right")
    ax1.grid(alpha=0.3)

    ax2.plot(fpr_test, tpr_test, color='darkgreen', lw=2, label=f'ROC (AUC = {auc_test:.2f})')
    ax2.plot([0, 1], [0, 1], color='navy', lw=2, linestyle='--')
    ax2.set_xlim([0.0, 1.0])
    ax2.set_ylim([0.0, 1.05])
    ax2.set_xlabel('False Positive Rate')
    ax2.set_ylabel('True Positive Rate')
    ax2.set_title('ROC Curve - Test Dataset')
    ax2.legend(loc="lower right")
    ax2.grid(alpha=0.3)

    plt.tight_layout()
    image_path = save_plot_to_file(fig, 'phase3_roc.png')
    if image_path:
        create_image_slide(prs, "Logistic Regression - ROC Curves", image_path,
                          "Model performance on training and test datasets")

    # Confusion Matrices
    from sklearn.metrics import confusion_matrix
    import seaborn as sns

    train_preds = (train_probs >= phase3.OPTIMAL_THRESHOLD)
    test_preds = (test_probs >= phase3.OPTIMAL_THRESHOLD)

    cm_train = confusion_matrix(y_train, train_preds)
    cm_test = confusion_matrix(y_test, test_preds)

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 5))

    sns.heatmap(cm_train, annot=True, fmt='d', cmap='Blues', ax=ax1,
                xticklabels=['Down', 'Up'], yticklabels=['Down', 'Up'])
    ax1.set_title(f'Confusion Matrix - Train\n(Threshold: {phase3.OPTIMAL_THRESHOLD})')
    ax1.set_xlabel('Predicted')
    ax1.set_ylabel('Actual')

    sns.heatmap(cm_test, annot=True, fmt='d', cmap='Greens', ax=ax2,
                xticklabels=['Down', 'Up'], yticklabels=['Down', 'Up'])
    ax2.set_title(f'Confusion Matrix - Test\n(Threshold: {phase3.OPTIMAL_THRESHOLD})')
    ax2.set_xlabel('Predicted')
    ax2.set_ylabel('Actual')

    plt.tight_layout()
    image_path = save_plot_to_file(fig, 'phase3_cm.png')
    if image_path:
        create_image_slide(prs, "Logistic Regression - Confusion Matrices", image_path,
                          "Classification performance with optimal threshold")    # ========================================================================
    # PHASE 4 - ML Models
    # ========================================================================

    create_title_slide(prs, "Phase 4", "Machine Learning Models Comparison")

    create_content_slide(prs, "Models Evaluated", [
        "1. Gaussian Naive Bayes",
        "   • Probabilistic classifier",
        "   • Fast training and prediction",
        "",
        "2. Decision Tree Classifier",
        "   • Rule-based decisions",
        "   • Easy to interpret",
        "",
        "3. Random Forest Classifier",
        "   • Ensemble of decision trees",
        "   • Robust and less prone to overfitting"
    ])

    # Train models
    print("Generating Phase 4 plots...")
    phase4.train_naive_bayes()
    phase4.train_decision_tree()
    phase4.train_random_forest()

    # Get the comparison early to populate AUC values
    phase4.plot_roc_curve_nb_train()
    plt.close()
    phase4.plot_roc_curve_nb_test()
    plt.close()
    phase4.plot_roc_curve_dt_train()
    plt.close()
    phase4.plot_roc_curve_dt_test()
    plt.close()
    phase4.plot_roc_curve_rf_train()
    plt.close()
    phase4.plot_roc_curve_rf_test()
    plt.close()

    # Now create ROC plots manually for presentation
    from sklearn.metrics import roc_curve, auc

    # Naive Bayes
    nb_train_probs = phase4.nb_model.predict_proba(phase4.X_train)[:, 1]
    nb_test_probs = phase4.nb_model.predict_proba(phase4.X_test)[:, 1]
    fpr_nb_train, tpr_nb_train, _ = roc_curve(phase4.y_train, nb_train_probs)
    fpr_nb_test, tpr_nb_test, _ = roc_curve(phase4.y_test, nb_test_probs)
    auc_nb_train = auc(fpr_nb_train, tpr_nb_train)
    auc_nb_test = auc(fpr_nb_test, tpr_nb_test)

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 5))
    ax1.plot(fpr_nb_train, tpr_nb_train, color='darkorange', lw=2, label=f'ROC (AUC = {auc_nb_train:.2f})')
    ax1.plot([0, 1], [0, 1], color='navy', lw=2, linestyle='--')
    ax1.set_title('Naive Bayes - Train Dataset')
    ax1.set_xlabel('False Positive Rate')
    ax1.set_ylabel('True Positive Rate')
    ax1.legend(loc="lower right")
    ax1.grid(alpha=0.3)

    ax2.plot(fpr_nb_test, tpr_nb_test, color='darkgreen', lw=2, label=f'ROC (AUC = {auc_nb_test:.2f})')
    ax2.plot([0, 1], [0, 1], color='navy', lw=2, linestyle='--')
    ax2.set_title('Naive Bayes - Test Dataset')
    ax2.set_xlabel('False Positive Rate')
    ax2.set_ylabel('True Positive Rate')
    ax2.legend(loc="lower right")
    ax2.grid(alpha=0.3)
    plt.tight_layout()

    image_path = save_plot_to_file(fig, 'phase4_nb_roc.png')
    if image_path:
        create_image_slide(prs, "Naive Bayes - ROC Curves", image_path)

    # Decision Tree
    dt_train_probs = phase4.dt_model.predict_proba(phase4.X_train)[:, 1]
    dt_test_probs = phase4.dt_model.predict_proba(phase4.X_test)[:, 1]
    fpr_dt_train, tpr_dt_train, _ = roc_curve(phase4.y_train, dt_train_probs)
    fpr_dt_test, tpr_dt_test, _ = roc_curve(phase4.y_test, dt_test_probs)
    auc_dt_train = auc(fpr_dt_train, tpr_dt_train)
    auc_dt_test = auc(fpr_dt_test, tpr_dt_test)

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 5))
    ax1.plot(fpr_dt_train, tpr_dt_train, color='darkorange', lw=2, label=f'ROC (AUC = {auc_dt_train:.2f})')
    ax1.plot([0, 1], [0, 1], color='navy', lw=2, linestyle='--')
    ax1.set_title('Decision Tree - Train Dataset')
    ax1.set_xlabel('False Positive Rate')
    ax1.set_ylabel('True Positive Rate')
    ax1.legend(loc="lower right")
    ax1.grid(alpha=0.3)

    ax2.plot(fpr_dt_test, tpr_dt_test, color='darkgreen', lw=2, label=f'ROC (AUC = {auc_dt_test:.2f})')
    ax2.plot([0, 1], [0, 1], color='navy', lw=2, linestyle='--')
    ax2.set_title('Decision Tree - Test Dataset')
    ax2.set_xlabel('False Positive Rate')
    ax2.set_ylabel('True Positive Rate')
    ax2.legend(loc="lower right")
    ax2.grid(alpha=0.3)
    plt.tight_layout()

    image_path = save_plot_to_file(fig, 'phase4_dt_roc.png')
    if image_path:
        create_image_slide(prs, "Decision Tree - ROC Curves", image_path)

    # Random Forest
    rf_train_probs = phase4.rf_model.predict_proba(phase4.X_train)[:, 1]
    rf_test_probs = phase4.rf_model.predict_proba(phase4.X_test)[:, 1]
    fpr_rf_train, tpr_rf_train, _ = roc_curve(phase4.y_train, rf_train_probs)
    fpr_rf_test, tpr_rf_test, _ = roc_curve(phase4.y_test, rf_test_probs)
    auc_rf_train = auc(fpr_rf_train, tpr_rf_train)
    auc_rf_test = auc(fpr_rf_test, tpr_rf_test)

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 5))
    ax1.plot(fpr_rf_train, tpr_rf_train, color='darkorange', lw=2, label=f'ROC (AUC = {auc_rf_train:.2f})')
    ax1.plot([0, 1], [0, 1], color='navy', lw=2, linestyle='--')
    ax1.set_title('Random Forest - Train Dataset')
    ax1.set_xlabel('False Positive Rate')
    ax1.set_ylabel('True Positive Rate')
    ax1.legend(loc="lower right")
    ax1.grid(alpha=0.3)

    ax2.plot(fpr_rf_test, tpr_rf_test, color='darkgreen', lw=2, label=f'ROC (AUC = {auc_rf_test:.2f})')
    ax2.plot([0, 1], [0, 1], color='navy', lw=2, linestyle='--')
    ax2.set_title('Random Forest - Test Dataset')
    ax2.set_xlabel('False Positive Rate')
    ax2.set_ylabel('True Positive Rate')
    ax2.legend(loc="lower right")
    ax2.grid(alpha=0.3)
    plt.tight_layout()

    image_path = save_plot_to_file(fig, 'phase4_rf_roc.png')
    if image_path:
        create_image_slide(prs, "Random Forest - ROC Curves", image_path)    # Model comparison would be a table - create as content slide
    comparison_df = phase4.compare_models()
    content = ["Model Performance Comparison (AUC Scores):"]
    for _, row in comparison_df.iterrows():
        content.append(f"{row['Model']}: Train={row['Train AUC']:.3f}, Test={row['Test AUC']:.3f}")
    create_content_slide(prs, "Model Comparison", content)

    # ========================================================================
    # PHASE 5 - Sentiment Analysis
    # ========================================================================

    create_title_slide(prs, "Phase 5", "Text Mining & Sentiment Analysis")

    create_content_slide(prs, "Text Analysis Approach", [
        "Data Source: Twitter/X posts about Nifty50",
        "Text Preprocessing Pipeline:",
        "  • Lowercase conversion",
        "  • Remove mentions, hashtags, URLs",
        "  • Remove punctuation and digits",
        "  • Remove stopwords",
        "  • Remove domain-specific words",
        "Analysis Methods:",
        "  • Word Cloud for topic identification",
        "  • VADER sentiment analysis"
    ])

    # Word Cloud
    print("Generating Phase 5 plots...")
    fig = phase5.generate_wordcloud()
    image_path = save_plot_to_file(fig, 'phase5_wordcloud.png')
    if image_path:
        create_image_slide(prs, "Word Cloud - Key Topics", image_path,
                          "Most frequent words in Nifty50 discussions")

    # Sentiment Pie
    fig = phase5.plot_sentiment_pie()
    image_path = save_plot_to_file(fig, 'phase5_sentiment_pie.png')
    if image_path:
        create_image_slide(prs, "Sentiment Breakdown", image_path,
                          "Proportion of positive, negative, and neutral sentiments")    # ========================================================================
    # CONCLUSIONS
    # ========================================================================

    create_title_slide(prs, "Conclusions", "Key Findings and Recommendations")

    create_content_slide(prs, "Key Findings", [
        "Global markets show strong correlations",
        "Significant predictors: HangSeng, Nikkei, DAX, VIX",
        "Machine learning models show comparable performance",
        "Random Forest slightly outperforms other models",
        "Sentiment analysis reveals market mood",
        "COVID-19 had significant impact on volatility"
    ])

    create_content_slide(prs, "Recommendations", [
        "Use ensemble of multiple models for predictions",
        "Monitor VIX for volatility signals",
        "Consider sentiment analysis as additional feature",
        "Regular model retraining with new data",
        "Combine technical and sentiment indicators",
        "Account for global market interdependencies"
    ])

    create_title_slide(prs, "Thank You", "Questions & Discussion")

    # Save presentation
    output_path = Path(__file__).parent / 'Market_Analysis_Presentation.pptx'
    prs.save(str(output_path))
    print(f"\n✅ Presentation saved to: {output_path}")

    return output_path


if __name__ == "__main__":
    print("Creating PowerPoint presentation from Market Analysis Dashboard...")
    print("This may take a few minutes as plots are being generated...\n")

    try:
        output_file = create_presentation()
        print("\n🎉 Success! Presentation created with all phases.")
        print(f"📊 File location: {output_file}")
    except Exception as e:
        print(f"\n❌ Error creating presentation: {e}")
        import traceback
        traceback.print_exc()
